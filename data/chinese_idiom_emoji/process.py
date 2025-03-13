import json
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import os



def extract_emoji_as_images(pptx_path, output_folder, gt_path):
    idioms = []
    prs = Presentation(pptx_path)
    json_data = []

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture = shape
                image_data = picture.image.blob
                image = Image.open(io.BytesIO(image_data))
                if image.mode == 'RGBA':
                    image = image.convert('RGB')
                image_name = f"c_idiom_emoji_{i+1}.jpeg"
                image.save(os.path.join(output_folder, image_name))
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if len(text) == 4:
                    idioms.append((i + 1, text))
                    json_data.append({
                        "image_name": image_name,
                        "gt": text
                    })
        print(f"==page {i+1}==")
        print(text)

    json_path = os.path.join(gt_path, "idioms.json")
    with open(json_path, 'w', encoding='utf-8') as json_file:
        json.dump(json_data, json_file, ensure_ascii=False, indent=4)

    return json_data

def get_image_size(image_path):
    with Image.open(image_path) as img:
        width, height = img.size
        print(f"pic width: {width}")
        print(f"pic height: {height}")
        return width, height

def crop_images_in_folder(folder_path, target_path, crop_width, crop_height, leftp, topp):
    for filename in os.listdir(folder_path):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
            file_path = os.path.join(folder_path, filename)
            cropped_img_path = os.path.join(target_path, filename)
            
            with Image.open(file_path) as img:
                width, height = img.size
                left = leftp
                top = topp
                right = left + crop_width
                bottom = top + crop_height
                
                cropped_img = img.crop((left, top, right, bottom))
                cropped_img.save(cropped_img_path)
                print(f"Cropped to: {cropped_img_path}")


if __name__ == "__main__":
    pptx_path = "/gemini/code/data/chinese_idiom_emoji/chinese_idiom_emoji.pptx"
    output_folder = "/gemini/code/data/chinese_idiom_emoji/img"
    gt_path = "/gemini/code/data/chinese_idiom_emoji"
    extract_emoji_as_images(pptx_path, output_folder, gt_path)

    # get_image_size("/gemini/code/data/chinese_idiom_emoji/img/c_idiom_emoji_125.jpeg")
    # folder_path = "/gemini/code/data/chinese_idiom_emoji/img"  
    # target_path = "/gemini/code/data/chinese_idiom_emoji/imgs"
    # crop_images_in_folder(folder_path, target_path, 2400, 800, 76, 350)
